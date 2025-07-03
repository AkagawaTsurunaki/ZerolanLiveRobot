import os.path

import PyPDF2
import docx2txt
from pptx import Presentation


def _pdf_to_text_pypdf2(path: str):
    """
    Extract text from PDF file.
    :param path:
    :return:
    """
    text = ""
    with open(path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
    return text


def _docx_to_text_simple(path: str):
    """
    Extract docx from DOCX file.
    :param path:
    :return:
    """
    return docx2txt.process(path)


def _pptx_to_text(path: str):
    """
    Extract pptx from PPTX file.
    :param path:
    :return:
    """
    prs = Presentation(path)
    text = ""

    for slide_num, slide in enumerate(prs.slides, 1):
        text += f"--- Slide {slide_num} ---\n"

        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

        text += "\n"

    return text


def extract_text(path: str):
    """
    Extract text from supported file.
    :param path:
    :return:
    """
    assert os.path.exists(path), f"File not found: {path}"
    suffix = path.split(".")[-1].lower()
    if suffix == "pdf":
        return _pdf_to_text_pypdf2(path)
    elif suffix == "docx":
        return _docx_to_text_simple(path)
    elif suffix == "pptx":
        return _pptx_to_text(path)
    else:
        raise NotImplementedError(f"Unsupported file type {suffix}")
